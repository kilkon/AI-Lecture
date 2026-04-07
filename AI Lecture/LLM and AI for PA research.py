# ============================================================
# 행정학 연구를 위한 AI LLM & AI Agent 활용 강의 PPT 자동 생성기
# v2.0 — Promptology 슬라이드 3장 추가 (총 30슬라이드)
# 실행환경: Google Colab 또는 Python 3.8+
# ============================================================

import subprocess, sys
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 1. 컬러 팔레트 ───────────────────────────────────────────
C_NAVY       = RGBColor(0x0D, 0x1B, 0x3E)
C_BLUE       = RGBColor(0x1A, 0x5C, 0xB5)
C_CYAN       = RGBColor(0x00, 0xC2, 0xD4)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
C_DARK_GRAY  = RGBColor(0x33, 0x33, 0x44)
C_GOLD       = RGBColor(0xF5, 0xA6, 0x23)
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)
C_RED        = RGBColor(0xE7, 0x4C, 0x3C)
C_PURPLE     = RGBColor(0x5B, 0x27, 0x9A)
C_TEAL       = RGBColor(0x00, 0x6A, 0x4E)

TOTAL_SLIDES = 30

# ── 2. 헬퍼 함수 ────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_name="맑은 고딕", font_size=18,
                bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def draw_header_bar(slide, title, subtitle=None,
                    bar_color=C_NAVY, accent_color=C_CYAN,
                    title_color=C_WHITE, sub_color=C_CYAN):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT_GRAY)
    add_rect(slide, 0, 0, 13.33, 1.4, fill_color=bar_color)
    add_rect(slide, 0, 1.4, 13.33, 0.05, fill_color=accent_color)
    add_rect(slide, 0, 0, 0.12, 1.4, fill_color=accent_color)
    add_textbox(slide, title, 0.25, 0.18, 12.5, 0.75,
                font_size=26, bold=True, color=title_color)
    if subtitle:
        add_textbox(slide, subtitle, 0.25, 0.88, 12.5, 0.45,
                    font_size=13, color=sub_color)

def draw_part_badge(slide, part_text, color=C_GOLD):
    add_rect(slide, 10.8, 0.12, 2.3, 0.45, fill_color=color)
    add_textbox(slide, part_text, 10.8, 0.12, 2.3, 0.45,
                font_size=12, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

def draw_slide_number(slide, num, total=TOTAL_SLIDES):
    add_textbox(slide, f"{num} / {total}", 11.8, 7.1, 1.3, 0.3,
                font_size=10, color=C_DARK_GRAY, align=PP_ALIGN.RIGHT)

def draw_bottom_bar(slide, color=C_NAVY):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_color=color)
    add_textbox(slide,
                "행정학 연구를 위한 AI LLM & AI Agent 활용  |  Promptology",
                0.2, 7.17, 10, 0.28,
                font_size=9, color=C_CYAN)

def info_box(slide, l, t, w, h, icon, title, body,
             box_color=C_NAVY, title_color=C_CYAN, body_color=C_WHITE):
    add_rect(slide, l, t, w, h, fill_color=box_color)
    add_rect(slide, l, t, w, 0.04, fill_color=C_CYAN)
    add_textbox(slide, icon, l+0.1, t+0.08, 0.6, 0.5,
                font_size=22, color=C_CYAN, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, l+0.65, t+0.1, w-0.75, 0.4,
                font_size=14, bold=True, color=title_color)
    add_textbox(slide, body, l+0.15, t+0.55, w-0.25, h-0.65,
                font_size=11, color=body_color, word_wrap=True)

# ── 3. 슬라이드 함수 ─────────────────────────────────────────

def slide_01_cover(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, 13.33, 7.5, fill_color=C_NAVY)
    add_rect(s, 8.5, 0, 4.83, 7.5, fill_color=RGBColor(0x0A,0x14,0x30))
    add_rect(s, 0, 2.8, 7.5, 0.06, fill_color=C_CYAN)
    add_rect(s, 0, 2.87, 7.5, 0.025, fill_color=C_BLUE)
    add_rect(s, 0.4, 0.45, 3.5, 0.42, fill_color=C_BLUE)
    add_textbox(s, "서울대 행정대학원 특강  |  2026", 0.4, 0.45, 3.5, 0.42,
                font_size=12, color=C_CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, "행정학 연구를 위한", 0.4, 1.05, 8.5, 0.75,
                font_size=30, color=C_LIGHT_GRAY)
    add_textbox(s, "AI LLM & AI Agent 활용", 0.4, 1.7, 8.5, 1.0,
                font_size=40, bold=True, color=C_WHITE)
    add_textbox(s,
                "NotebookLM · Genspark · Promptology · Vibe Coding으로\n바꾸는 연구 패러다임",
                0.4, 2.98, 8.2, 1.1, font_size=17, color=C_CYAN)
    add_rect(s, 0.4, 4.25, 3.0, 0.04, fill_color=C_GOLD)
    add_textbox(s, "행정학과 대학원  |  연구방법론 특강",
                0.4, 4.42, 7.0, 0.4, font_size=13, color=C_LIGHT_GRAY)
    for i, txt in enumerate(["Research","Question","Promptology","Analysis","Insight"]):
        fc = C_WHITE if i == 2 else RGBColor(0x2A,0x3F,0x6F)
        sz = 24 if i == 2 else 20
        add_textbox(s, txt, 9.0, 1.2+i*1.0, 3.5, 0.6,
                    font_size=sz, bold=(i==2), color=fc,
                    align=PP_ALIGN.CENTER)

def slide_02_roadmap(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "강의 로드맵",
                    "오늘 1시간 동안 배우는 것 — 연구 전 과정을 AI와 함께")
    draw_bottom_bar(s)
    draw_slide_number(s, 2)

    steps = [
        ("🔍","연구질문\n설계","Genspark"),
        ("📚","선행연구\n검토","NotebookLM"),
        ("✍️","프롬프톨로지\n(Promptology)","SPARRO/CRAFT"),
        ("🏗️","연구모형\n설계","Claude 활용"),
        ("💻","통계분석\n실행","Vibe Coding"),
    ]
    box_w, box_h, start_l, top = 2.1, 3.2, 0.45, 1.9
    for i,(icon,title,tool) in enumerate(steps):
        l = start_l + i*2.55
        col = C_BLUE if i%2==0 else RGBColor(0x0D,0x47,0xA1)
        if i == 2:
            col = C_PURPLE
        add_rect(s, l, top, box_w, box_h, fill_color=col)
        add_rect(s, l, top, box_w, 0.06, fill_color=C_CYAN if i!=2 else C_GOLD)
        add_textbox(s, icon, l, top+0.12, box_w, 0.55,
                    font_size=24, align=PP_ALIGN.CENTER)
        add_textbox(s, title, l, top+0.65, box_w, 0.9,
                    font_size=14, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_rect(s, l+0.3, top+1.6, box_w-0.6, 0.04, fill_color=C_CYAN)
        add_textbox(s, tool, l, top+1.72, box_w, 0.9,
                    font_size=11,
                    color=C_GOLD if i==2 else C_CYAN,
                    align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "▶", l+box_w+0.1, top+1.2, 0.35, 0.5,
                        font_size=14, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

    add_rect(s, 0.45, 5.35, 12.43, 0.62,
             fill_color=RGBColor(0x0D,0x1B,0x3E))
    add_textbox(s,
                '💡 핵심: "좋은 프롬프트가 좋은 연구를 만든다" — Promptology는 AI 시대의 연구 방법론이다',
                0.6, 5.42, 12.1, 0.45,
                font_size=14, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

def slide_03_why(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "왜 행정학 연구자가 AI를 배워야 하는가?",
                    "연구 대상(공공행정)이 AI로 변하고 있는데, 연구 방법은 그대로?")
    draw_part_badge(s, "PART 1 | 기초 이해")
    draw_bottom_bar(s)
    draw_slide_number(s, 3)

    stats = [
        ("💰","AI Agent 시장","2024년 54억 달러\n→ 2030년 470억 달러\nCAGR 45.8%"),
        ("🏛️","공공부문 도입","기업 40% 이상이\nAI Agent를\n워크플로우 통합"),
        ("📉","연구의 위기","행정 현장은 AI화\n연구 방법론은\n여전히 전통 방식"),
    ]
    for i,(icon,title,body) in enumerate(stats):
        info_box(s, 0.4+i*4.3, 1.65, 3.9, 2.85, icon, title, body)

    add_rect(s, 0.4, 4.75, 12.53, 0.9, fill_color=C_RED)
    add_textbox(s,
                '❓ "AI를 모르는 행정학 연구자가 AI 행정을 연구할 수 있는가?"',
                0.6, 4.85, 12.1, 0.65,
                font_size=16, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s,
                "→ 20년 전 행정학자에게 SPSS는 낯선 도구였다. 지금은 기본 소양이다. LLM도 마찬가지다.",
                0.4, 5.85, 12.53, 0.45,
                font_size=13, color=C_DARK_GRAY,
                align=PP_ALIGN.CENTER)

def slide_04_llm_vs_agent(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "LLM과 AI Agent: 무엇이 어떻게 다른가?",
                    "단순 질답 도구 vs 목표를 스스로 달성하는 시스템")
    draw_part_badge(s, "PART 1 | 기초 이해")
    draw_bottom_bar(s)
    draw_slide_number(s, 4)

    add_rect(s, 0.4, 1.65, 5.9, 4.85, fill_color=C_BLUE)
    add_rect(s, 0.4, 1.65, 5.9, 0.06, fill_color=C_CYAN)
    add_textbox(s, "🤖  LLM", 0.4, 1.72, 5.9, 0.6,
                font_size=22, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "예: ChatGPT · Claude · Gemini",
                0.4, 2.28, 5.9, 0.38, font_size=13, color=C_CYAN,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(["질문에 답하는 도구",
                              "입력 → 출력 (1회성)",
                              "텍스트 분석 · 요약 · 코딩 보조",
                              "언어 이해 · 생성"]):
        add_rect(s, 0.6, 2.85+i*0.68, 5.5, 0.55,
                 fill_color=RGBColor(0x0D,0x47,0xA1))
        add_textbox(s, f"✦  {item}", 0.75, 2.88+i*0.68, 5.2, 0.5,
                    font_size=13, color=C_WHITE)

    add_rect(s, 7.0, 1.65, 5.9, 4.85, fill_color=C_TEAL)
    add_rect(s, 7.0, 1.65, 5.9, 0.06, fill_color=C_GOLD)
    add_textbox(s, "🚀  AI Agent", 7.0, 1.72, 5.9, 0.6,
                font_size=22, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "예: Genspark · AutoGen · CrewAI",
                7.0, 2.28, 5.9, 0.38, font_size=13, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(["목표를 스스로 달성하는 시스템",
                              "계획 → 실행 → 검증 → 반복",
                              "연구 파이프라인 전체 자동화",
                              "도구사용 + 기억 + 자기수정"]):
        add_rect(s, 7.2, 2.85+i*0.68, 5.5, 0.55,
                 fill_color=RGBColor(0x00,0x4D,0x38))
        add_textbox(s, f"✦  {item}", 7.35, 2.88+i*0.68, 5.2, 0.5,
                    font_size=13, color=C_WHITE)

    add_rect(s, 6.1, 2.8, 0.8, 0.8, fill_color=C_GOLD)
    add_textbox(s, "VS", 6.1, 2.82, 0.8, 0.75,
                font_size=20, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

def slide_05_tools(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "오늘 강의에서 활용할 AI 도구 5종",
                    "각 도구는 연구의 특정 단계에 최적화 — Promptology가 이 모든 도구를 연결한다")
    draw_part_badge(s, "PART 1 | 기초 이해")
    draw_bottom_bar(s)
    draw_slide_number(s, 5)

    tools = [
        ("🔍","Genspark","AI Super Agent",
         "연구질문 탐색 · 딥리서치\n슬라이드 자동 생성까지\n올인원 연구 플랫폼",
         C_BLUE),
        ("📚","NotebookLM","Google AI 문헌 분석",
         "업로드 소스만 참조하는\n신뢰 기반 문헌 분석 도구\n할루시네이션 위험 최소화",
         C_TEAL),
        ("✍️","Promptology","프롬프트 설계 학문",
         "AI와의 상호작용을 최적화하는\n전략적 프롬프트 설계 원칙\nSPARRO·CRAFT 프레임워크",
         C_PURPLE),
        ("🤖","Claude / GPT","대화형 연구 파트너",
         "연구모형 설계 · 가설 작성\n이론 탐색 · 설문 문항 생성\n프롬프톨로지 기법 적용",
         RGBColor(0x8B,0x45,0x13)),
        ("💻","Cursor / Colab","바이브 코딩 환경",
         "자연어로 Python 코드 생성\n통계분석 자동 실행\n결과 해석 초안 생성",
         RGBColor(0x8B,0x00,0x00)),
    ]
    # 상단 3개
    for i,(icon,name,sub,body,col) in enumerate(tools[:3]):
        l = 0.4 + i*4.28
        add_rect(s, l, 1.65, 3.98, 2.3, fill_color=col)
        add_rect(s, l, 1.65, 3.98, 0.05,
                 fill_color=C_GOLD if i==2 else C_CYAN)
        add_textbox(s, icon, l+0.1, t:=1.75, 0.7, 0.6,
                    font_size=24, color=C_WHITE)
        add_textbox(s, name, l+0.75, 1.75, 3.0, 0.45,
                    font_size=16, bold=True, color=C_WHITE)
        add_textbox(s, sub,  l+0.75, 2.18, 3.0, 0.35,
                    font_size=11,
                    color=C_GOLD if i==2 else C_CYAN)
        add_textbox(s, body, l+0.15, 2.6, 3.65, 1.2,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)
    # 하단 2개
    for i,(icon,name,sub,body,col) in enumerate(tools[3:]):
        l = 1.77 + i*5.3
        add_rect(s, l, 4.18, 4.8, 2.3, fill_color=col)
        add_rect(s, l, 4.18, 4.8, 0.05, fill_color=C_CYAN)
        add_textbox(s, icon, l+0.1, 4.28, 0.7, 0.6,
                    font_size=24, color=C_WHITE)
        add_textbox(s, name, l+0.75, 4.28, 3.8, 0.45,
                    font_size=16, bold=True, color=C_WHITE)
        add_textbox(s, sub,  l+0.75, 4.7, 3.8, 0.35,
                    font_size=11, color=C_CYAN)
        add_textbox(s, body, l+0.15, 5.12, 4.5, 1.2,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)

def slide_06_principles(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "AI 활용 연구의 핵심 원칙 3가지",
                    "도구를 사용하기 전에 — 방법론적·윤리적 기반을 먼저")
    draw_part_badge(s, "PART 1 | 기초 이해")
    draw_bottom_bar(s)
    draw_slide_number(s, 6)

    principles = [
        ("①","Human-in-the-Loop",
         "AI는 보조, 최종 판단은 연구자",
         "AI가 분석 초안을 만들어도 연구자가 반드시 검토·수정·판단해야 한다. "
         "연구의 지적 책임은 항상 인간 연구자에게 있다.",C_BLUE),
        ("②","출처 검증 필수 — 할루시네이션 경계",
         "AI는 그럴듯한 거짓말을 한다",
         "LLM은 존재하지 않는 논문, 틀린 통계를 자신있게 생성할 수 있다. "
         "모든 AI 출력은 원출처와 대조 검증이 필수이다.",C_RED),
        ("③","투명한 보고",
         "논문에 AI 활용 방식을 명시해야 한다",
         "어떤 AI 도구를 어느 연구 단계에 활용했는지 방법론 섹션에 기술하는 것이 "
         "국제 학술 규범(COREQ+LLM, 2025)으로 자리잡고 있다.",C_TEAL),
    ]
    for i,(num,title,sub,body,col) in enumerate(principles):
        t = 1.65 + i*1.72
        add_rect(s, 0.4, t, 12.53, 1.5, fill_color=col)
        add_rect(s, 0.4, t, 0.65, 1.5, fill_color=RGBColor(0,0,0))
        add_textbox(s, num, 0.4, t, 0.65, 1.5,
                    font_size=24, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title, 1.15, t+0.1, 5.5, 0.45,
                    font_size=15, bold=True, color=C_WHITE)
        add_textbox(s, sub,   1.15, t+0.52, 5.0, 0.35,
                    font_size=12,
                    color=C_GOLD if col==C_RED else C_CYAN)
        add_textbox(s, body,  6.5,  t+0.12, 6.2, 1.2,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)

# ══════════════════════════════════════════════════════════════
# ★★★ NEW SLIDE 7 — Promptology 개념 정의 ★★★
# ══════════════════════════════════════════════════════════════
def slide_07_promptology_def(prs):
    """Promptology 개념 정의 슬라이드"""
    s = blank_slide(prs)
    draw_header_bar(s,
                    "Promptology — AI 시대의 새로운 연구 방법론 언어",
                    "MDPI Information, 2024 | '프롬프트를 연구하는 학문' — Olla et al.")
    draw_part_badge(s, "★ Promptology", color=C_PURPLE)
    draw_bottom_bar(s)
    draw_slide_number(s, 7)

    # 정의 박스
    add_rect(s, 0.4, 1.65, 12.53, 1.1, fill_color=C_PURPLE)
    add_rect(s, 0.4, 1.65, 0.1, 1.1, fill_color=C_GOLD)
    add_textbox(s,
                '"Promptology is an interdisciplinary field that focuses on designing '
                'strategic and secure prompts for generative AI systems. It integrates '
                'technical skills, cognitive science, and frameworks to optimize '
                'human-AI interactions."',
                0.65, 1.75, 12.0, 0.85,
                font_size=13, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "— Olla et al., MDPI Information, 2024",
                9.5, 2.7, 3.2, 0.35,
                font_size=10, color=C_GOLD)

    # 3개 키 개념
    concepts = [
        ("🔬","학문적 정의",
         "단순한 기술 스킬이 아닌\n독립적 학제간 연구 분야\n(언어학 + 인지과학 + AI)"),
        ("🎯","핵심 목적",
         "AI와의 상호작용을\n최대한 효율적·윤리적·\n정확하게 최적화"),
        ("🔗","연구와의 연결",
         "프롬프트 설계 능력 =\n연구 설계 능력\nAI 시대의 연구 방법론"),
    ]
    for i,(icon,title,body) in enumerate(concepts):
        l = 0.4 + i*4.3
        add_rect(s, l, 3.1, 3.9, 2.65, fill_color=C_NAVY)
        add_rect(s, l, 3.1, 3.9, 0.06, fill_color=C_PURPLE)
        add_textbox(s, icon, l, 3.2, 3.9, 0.6,
                    font_size=26, align=PP_ALIGN.CENTER)
        add_textbox(s, title, l, 3.78, 3.9, 0.45,
                    font_size=14, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, body, l+0.15, 4.28, 3.6, 1.2,
                    font_size=12, color=C_CYAN,
                    align=PP_ALIGN.CENTER, word_wrap=True)

    # 비유 박스
    add_rect(s, 0.4, 5.95, 12.53, 0.72, fill_color=C_NAVY)
    add_textbox(s,
                '💡 비유:  언어학이 인간 언어를 연구하듯, Promptology는 "AI와 소통하는 언어"를 연구한다',
                0.6, 6.05, 12.1, 0.52,
                font_size=13, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# ★★★ NEW SLIDE 8 — Promptology 핵심 기법 ★★★
# ══════════════════════════════════════════════════════════════
def slide_08_promptology_techniques(prs):
    """Promptology 핵심 프롬프팅 기법"""
    s = blank_slide(prs)
    draw_header_bar(s,
                    "Promptology 핵심 기법 — 행정학 연구에 바로 적용하는 6가지 전략",
                    "Zero-shot부터 Chain-of-Thought까지 — 연구 단계별 최적 기법 선택")
    draw_part_badge(s, "★ Promptology", color=C_PURPLE)
    draw_bottom_bar(s)
    draw_slide_number(s, 8)

    techniques = [
        ("① Zero-shot",
         "예시 없이 직접 질문",
         '"이 정책 문서의 핵심 논거를\n3가지로 요약해줘"',
         "간단한 요약·분류 작업에 적합\n빠르지만 정확도 편차 있음",
         C_BLUE),
        ("② Few-shot",
         "예시를 제공하여 패턴 학습",
         '"아래 예시처럼 분석해줘:\n[예시1]→[출력1]...\n이제 이 문서를 분석해줘"',
         "일관된 형식의 반복 작업에 강력\n코딩·분류·평가에 효과적",
         C_TEAL),
        ("③ Chain-of-Thought",
         "단계별 사고 과정 유도",
         '"단계별로 생각해서\n이 정책의 인과관계를\n분석해줘"',
         "복잡한 논리 추론에 최적\n연구모형 설계·가설 검토",
         C_PURPLE),
        ("④ Role Prompting",
         "전문가 역할 부여",
         '"너는 행정학 연구방법론\n전문가야. 다음 가설의\n타당성을 검토해줘"',
         "도메인 전문성 활용에 탁월\n연구설계·논문 리뷰에 활용",
         RGBColor(0x8B,0x45,0x13)),
        ("⑤ Structured Output",
         "출력 형식 명시 지정",
         '"결과를 JSON 형식으로\n출력해줘: {저자, 연도,\n주요발견, 방법론}"',
         "데이터 추출·표 작성에 필수\nNotebookLM 활용 시 강점",
         RGBColor(0x00,0x5C,0x8A)),
        ("⑥ Iterative Refinement",
         "반복적 개선 프롬프팅",
         '"앞의 답변에서 행정학\n이론과의 연결이 부족해.\n이 부분을 보완해줘"',
         "연구 초안 개선에 가장 실용적\n한 번으로 끝내려 하지 말 것",
         RGBColor(0x7B,0x1F,0x1F)),
    ]
    for i,(name,desc,example,tip,col) in enumerate(techniques):
        l = 0.35 + (i%3)*4.32
        t = 1.62 + (i//3)*2.72
        add_rect(s, l, t, 4.05, 2.55, fill_color=col)
        add_rect(s, l, t, 4.05, 0.05, fill_color=C_GOLD if i%2==0 else C_CYAN)
        add_textbox(s, name, l+0.1, t+0.1, 3.8, 0.42,
                    font_size=14, bold=True, color=C_WHITE)
        add_textbox(s, desc, l+0.1, t+0.5, 3.8, 0.35,
                    font_size=11, color=C_CYAN)
        add_rect(s, l+0.1, t+0.9, 3.75, 0.85,
                 fill_color=RGBColor(0,0,0))
        add_textbox(s, example, l+0.18, t+0.93, 3.55, 0.78,
                    font_name="Courier New", font_size=9,
                    color=RGBColor(0x00,0xFF,0x7F), word_wrap=True)
        add_textbox(s, tip, l+0.1, t+1.82, 3.75, 0.62,
                    font_size=10, color=C_LIGHT_GRAY, word_wrap=True)

# ══════════════════════════════════════════════════════════════
# ★★★ NEW SLIDE 9 — SPARRO 프레임워크 ★★★
# ══════════════════════════════════════════════════════════════
def slide_09_sparro(prs):
    """SPARRO 프레임워크 — 학술 연구에서의 AI 활용 구조"""
    s = blank_slide(prs)
    draw_header_bar(s,
                    "SPARRO 프레임워크 — 학술 연구에서 AI를 윤리적으로 활용하는 6단계",
                    "Olla et al. (2024) | Promptology 기반 AI 연구 통합 가이드라인")
    draw_part_badge(s, "★ Promptology", color=C_PURPLE)
    draw_bottom_bar(s)
    draw_slide_number(s, 9)

    sparro = [
        ("S","Strategy\n전략 수립",
         "AI 활용 범위·역할·한계를\n사전에 명확히 계획\n→ AI 활용 선언문 작성",
         C_BLUE),
        ("P","Prompt Design\n프롬프트 설계",
         "CRAFT 모델 적용\n(명확성·근거·대상·형식·과제)\n→ 고품질 출력 확보",
         C_PURPLE),
        ("A","Adopting\n출력 통합",
         "AI 생성 내용을 연구자\n자신의 논문에 통합\n→ 목소리·논지 유지",
         C_TEAL),
        ("R","Reviewing\n내용 검토",
         "정확성·관련성·일관성\n비판적 평가 + 원출처 대조\n→ 할루시네이션 제거",
         RGBColor(0x8B,0x45,0x13)),
        ("R","Refining\n반복 개선",
         "논거 재구성·언어 다듬기\n연구자 통찰 추가\n→ 학술적 수준 확보",
         RGBColor(0x7B,0x1F,0x1F)),
        ("O","Optimizing\n최적화",
         "표절 검사 + 출처 검증\n원본성·학술 정직성 확인\n→ 최종 논문 완성",
         RGBColor(0x1A,0x5C,0x2A)),
    ]
    box_w, box_h = 2.05, 3.8
    for i,(letter,title,body,col) in enumerate(sparro):
        l = 0.38 + i*2.09
        add_rect(s, l, 1.62, box_w, box_h, fill_color=col)
        add_rect(s, l, 1.62, box_w, 0.06,
                 fill_color=C_GOLD if i%2==0 else C_CYAN)
        # 대문자
        add_rect(s, l, 1.62, box_w, 0.75,
                 fill_color=RGBColor(0,0,0))
        add_textbox(s, letter, l, 1.62, box_w, 0.75,
                    font_size=32, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        # 단계명
        add_textbox(s, title, l+0.05, 2.42, box_w-0.1, 0.72,
                    font_size=12, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        # 내용
        add_textbox(s, body, l+0.08, 3.18, box_w-0.16, 1.98,
                    font_size=10, color=C_LIGHT_GRAY,
                    align=PP_ALIGN.CENTER, word_wrap=True)
        # 화살표
        if i < 5:
            add_textbox(s, "▶", l+box_w+0.01, 3.3, 0.1, 0.4,
                        font_size=9, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

    # CRAFT 모델 설명 박스
    add_rect(s, 0.38, 5.6, 12.57, 1.08, fill_color=C_NAVY)
    add_rect(s, 0.38, 5.6, 12.57, 0.05, fill_color=C_PURPLE)
    add_textbox(s, "📐  CRAFT 모델 (Prompt Design 세부 원칙)",
                0.58, 5.65, 6.0, 0.38,
                font_size=12, bold=True, color=C_PURPLE)
    craft = [("C","Clarity\n명확성"),("R","Rationale\n근거"),
             ("A","Audience\n대상"),("F","Format\n형식"),
             ("T","Tasks\n과제")]
    for i,(letter,desc) in enumerate(craft):
        l = 0.58 + i*2.45
        add_rect(s, l, 6.1, 2.2, 0.48,
                 fill_color=C_PURPLE if i%2==0 else RGBColor(0x3A,0x0A,0x6A))
        add_textbox(s, f"{letter} — {desc}", l+0.08, 6.13, 2.0, 0.4,
                    font_size=11, bold=True, color=C_WHITE)

def slide_10_rq_design(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "연구 질문(RQ) 설계 — 왜 어렵고, AI가 어떻게 도울 수 있는가?",
                    "대학원생이 가장 오래 막히는 단계를 AI가 가속화한다")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 10)

    add_rect(s, 0.4, 1.65, 5.8, 4.5, fill_color=RGBColor(0x6B,0x6B,0x7B))
    add_rect(s, 0.4, 1.65, 5.8, 0.5, fill_color=RGBColor(0x4A,0x4A,0x5A))
    add_textbox(s, "😓  기존 방식", 0.4, 1.68, 5.8, 0.45,
                font_size=16, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,st in enumerate(["관심 분야 막연히 설정",
                            "교수 상담 → 피드백 대기",
                            "논문 수십 편 수작업 검색",
                            "연구 공백 파악에 수 주 소요",
                            "RQ 초안 → 다시 수정 반복"]):
        add_textbox(s, f"  {i+1}.  {st}", 0.6, 2.3+i*0.7, 5.3, 0.55,
                    font_size=13, color=C_LIGHT_GRAY)

    add_rect(s, 7.1, 1.65, 5.8, 4.5, fill_color=C_BLUE)
    add_rect(s, 7.1, 1.65, 5.8, 0.5, fill_color=RGBColor(0x0D,0x3A,0x8A))
    add_textbox(s, "🚀  AI + Promptology 방식", 7.1, 1.68, 5.8, 0.45,
                font_size=16, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,st in enumerate(["Role 프롬프트로 전문가 모드 설정",
                            "Genspark Deep Research 실행",
                            "연구 공백(Gap) AI가 자동 도출",
                            "Chain-of-Thought로 RQ 정교화",
                            "FINER 기준 자동 평가 요청"]):
        add_textbox(s, f"  {i+1}.  {st}", 7.3, 2.3+i*0.7, 5.3, 0.55,
                    font_size=13, color=C_WHITE)

    add_textbox(s, "→", 6.3, 3.55, 0.7, 0.7,
                font_size=30, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

def slide_11_genspark_demo(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "🔧 [실습 데모] Genspark으로 연구 질문 설계하기",
                    "Deep Research Agent + Promptology 기법 적용")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 11)

    add_rect(s, 0.4, 1.65, 12.53, 1.65, fill_color=RGBColor(0x1A,0x1A,0x2E))
    add_rect(s, 0.4, 1.65, 0.08, 1.65, fill_color=C_PURPLE)
    add_textbox(s, "💬  Promptology 적용 — Role + Chain-of-Thought 기법",
                0.6, 1.7, 12.0, 0.38,
                font_size=12, bold=True, color=C_PURPLE)
    add_textbox(s,
                "[역할] 너는 행정학 연구방법론 전문가야.\n"
                "[단계적 사고] 다음 순서로 분석해줘:\n"
                "① 한국 지방정부 디지털 전환·주민 신뢰 연구의 현황 정리 →\n"
                "② 연구 공백(research gap) 3가지 도출 →\n"
                "③ FINER 기준을 충족하는 연구 질문(RQ) 5개 제안",
                0.6, 2.07, 12.0, 1.1,
                font_size=12, color=C_WHITE, word_wrap=True)

    steps = [
        ("🌐","실시간\n웹 탐색","학술 DB\n최신 논문"),
        ("📊","연구 흐름\n시각화","트렌드\n자동 정리"),
        ("🔭","Gap\n분석","미개척 영역\n도출"),
        ("❓","RQ 초안\n생성","5개 즉시\n생성"),
        ("✅","FINER\n자동 평가","기준 충족\n검토"),
    ]
    for i,(icon,title,body) in enumerate(steps):
        l = 0.4 + i*2.5
        add_rect(s, l, 3.55, 2.2, 2.75, fill_color=C_BLUE)
        add_rect(s, l, 3.55, 2.2, 0.05, fill_color=C_CYAN)
        add_textbox(s, icon, l, 3.65, 2.2, 0.55,
                    font_size=22, align=PP_ALIGN.CENTER)
        add_textbox(s, title, l, 4.18, 2.2, 0.65,
                    font_size=13, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, body, l+0.1, 4.82, 2.0, 0.88,
                    font_size=11, color=C_CYAN,
                    align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "▶", l+2.2, 4.7, 0.3, 0.45,
                        font_size=12, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

def slide_12_finer(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "좋은 연구 질문의 조건 — FINER 기준 + AI 정교화",
                    "Promptology 기법: Iterative Refinement로 RQ를 다듬는다")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 12)

    finer = [
        ("F","Feasible","실현 가능한가?","주어진 자원·시간 안에 수행 가능한가"),
        ("I","Interesting","흥미로운가?","학계와 실무에 의미 있는 기여인가"),
        ("N","Novel","새로운가?","기존 연구에서 다루지 않은 측면인가"),
        ("E","Ethical","윤리적인가?","연구 대상·데이터 활용에 문제없는가"),
        ("R","Relevant","적합한가?","행정학 이론·정책 실무와 연결되는가"),
    ]
    for i,(letter,eng,kor,desc) in enumerate(finer):
        t = 1.65 + i*1.05
        add_rect(s, 0.4, t, 0.7, 0.85, fill_color=C_CYAN)
        add_textbox(s, letter, 0.4, t, 0.7, 0.85,
                    font_size=26, bold=True, color=C_NAVY,
                    align=PP_ALIGN.CENTER)
        add_rect(s, 1.15, t, 11.78, 0.85,
                 fill_color=C_NAVY if i%2==0 else C_BLUE)
        add_textbox(s, f"{eng}  |  {kor}", 1.3, t+0.05, 5.0, 0.42,
                    font_size=14, bold=True, color=C_WHITE)
        add_textbox(s, desc, 6.5, t+0.08, 6.2, 0.65,
                    font_size=12, color=C_CYAN)

    add_rect(s, 0.4, 7.0, 12.53, 0.52, fill_color=C_PURPLE)
    add_textbox(s,
                '✍️ Promptology 적용:  "이 RQ가 FINER 기준을 충족하는지 평가하고, '
                'Chain-of-Thought로 개선안을 단계별로 제시해줘"',
                0.6, 7.05, 12.1, 0.4,
                font_size=11, bold=True, color=C_WHITE)

def slide_13_notebooklm(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "NotebookLM — 내가 업로드한 문헌만 참조하는 AI",
                    '"내 논문 더미를 읽어주는 AI 조교" — 할루시네이션 위험 최소화')
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 13)

    add_rect(s, 0.4, 1.65, 12.53, 0.8, fill_color=C_TEAL)
    add_textbox(s,
                "🔑 핵심: 인터넷 전체가 아닌, 연구자가 직접 업로드한 소스만 참조 → 인용 추적 가능 + 신뢰성 높음",
                0.6, 1.72, 12.1, 0.6,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(s, "📂  지원 소스", 0.4, 2.65, 6.0, 0.42,
                font_size=14, bold=True, color=C_NAVY)
    for i,src in enumerate(["📄 PDF 논문·보고서","📝 Google Docs",
                             "🌐 웹 URL","🎥 YouTube 영상",
                             "🎙️ 오디오 파일","📊 텍스트 파일"]):
        l = 0.4+(i%3)*2.0; t = 3.1+(i//3)*0.62
        add_rect(s, l, t, 1.85, 0.5, fill_color=C_BLUE)
        add_textbox(s, src, l+0.05, t+0.05, 1.75, 0.4,
                    font_size=11, color=C_WHITE)

    add_textbox(s, "⚙️  주요 기능 + Promptology 활용법", 7.0, 2.65, 6.0, 0.42,
                font_size=14, bold=True, color=C_NAVY)
    for i,func in enumerate([
        "📋 자동 요약 (Structured Output 프롬프트)",
        "🔗 인용 출처 추적 + 소스 비교",
        "🎙️ 오디오 오버뷰 자동 생성",
        "🔬 심층 연구(Deep Research) 모드",
        "❓ Few-shot 기반 문헌 분류 질문",
    ]):
        add_rect(s, 7.0, 3.1+i*0.62, 6.28, 0.5,
                 fill_color=RGBColor(0x00,0x4D,0x38) if i%2==0 else C_TEAL)
        add_textbox(s, func, 7.15, 3.13+i*0.62, 6.0, 0.42,
                    font_size=12, color=C_WHITE)

def slide_14_notebooklm_demo(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "🔧 [실습 데모] NotebookLM으로 선행연구 검토하기",
                    "Structured Output + Few-shot 프롬프트로 선행연구 매트릭스 자동 완성")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 14)

    steps_data = [
        ("STEP 1","소스 업로드",
         "관련 논문 10~20편 PDF\nRISS·KISS 한국 논문\n기존 보유 자료 전체",
         C_BLUE),
        ("STEP 2","Structured Output 질문",
         '"저자·연도·이론·방법론·\n주요발견을 JSON 표로\n출력해줘 (Few-shot)"',
         C_TEAL),
        ("STEP 3","매트릭스 생성",
         '"공통 이론·매개변수·\n연구 방법론 분류"\n→ 선행연구 표 자동 완성',
         C_PURPLE),
    ]
    for i,(step,title,body,col) in enumerate(steps_data):
        l = 0.4+i*4.3
        add_rect(s, l, 1.65, 3.9, 4.1, fill_color=col)
        add_rect(s, l, 1.65, 3.9, 0.06,
                 fill_color=C_CYAN if i<2 else C_GOLD)
        add_rect(s, l, 1.65, 3.9, 0.52,
                 fill_color=RGBColor(0,0,0))
        add_textbox(s, step, l, 1.68, 3.9, 0.45,
                    font_size=14, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title, l, 2.25, 3.9, 0.52,
                    font_size=16, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, body, l+0.15, 2.88, 3.6, 2.5,
                    font_size=12, color=C_LIGHT_GRAY, word_wrap=True)
        if i < 2:
            add_textbox(s, "▶", l+3.9+0.15, 3.5, 0.3, 0.5,
                        font_size=18, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 5.95, 12.53, 0.85, fill_color=RGBColor(0x1A,0x1A,0x2E))
    add_textbox(s,
                "📊 출력: 저자 | 연도 | 이론 | 방법론 | 주요 발견  →  선행연구 매트릭스 자동 완성\n"
                "✍️ Promptology 포인트: Structured Output 기법으로 일관된 형식 강제",
                0.6, 6.02, 12.1, 0.7,
                font_size=12, color=C_CYAN,
                align=PP_ALIGN.CENTER)

def slide_15_methodology(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "NotebookLM 활용의 방법론적 의미와 한계",
                    "AI는 패턴을 인식한다 — 연구자의 해석적 판단이 여전히 핵심")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 15)

    add_rect(s, 0.4, 1.65, 5.8, 3.2, fill_color=RGBColor(0x5A,0x5A,0x6A))
    add_textbox(s, "🕰️  전통적 방식", 0.4, 1.68, 5.8, 0.48,
                font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for txt in ["각 논문 직접 정독","수작업 메모·분류",
                "선행연구 표 수동 작성","소요: 수 주 ~ 수 개월"]:
        idx = ["각 논문 직접 정독","수작업 메모·분류",
               "선행연구 표 수동 작성","소요: 수 주 ~ 수 개월"].index(txt)
        add_textbox(s, f"  •  {txt}", 0.55, 2.28+idx*0.52, 5.4, 0.45,
                    font_size=13, color=C_LIGHT_GRAY)

    add_rect(s, 7.1, 1.65, 5.8, 3.2, fill_color=C_TEAL)
    add_textbox(s, "🚀  NotebookLM + Promptology", 7.1, 1.68, 5.8, 0.48,
                font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for txt in ["업로드 → AI 분류·요약 자동화",
                "Structured Output으로 표 즉시 생성",
                "Iterative Refinement로 분석 심화",
                "소요: 수 시간"]:
        idx = ["업로드 → AI 분류·요약 자동화",
               "Structured Output으로 표 즉시 생성",
               "Iterative Refinement로 분석 심화",
               "소요: 수 시간"].index(txt)
        add_textbox(s, f"  •  {txt}", 7.25, 2.28+idx*0.52, 5.4, 0.45,
                    font_size=13, color=C_WHITE)

    add_textbox(s, "⚡", 6.3, 2.9, 0.7, 0.6,
                font_size=28, align=PP_ALIGN.CENTER)

    warns = [
        ("⚠️","이해 vs 패턴 인식",
         "AI는 논문을 '이해'하는 것이 아니라 '패턴을 인식'한다. 연구자의 해석적 판단 필수"),
        ("📋","COREQ+LLM (2025)",
         "질적 연구에서 LLM 활용 보고 국제 표준. SPARRO 框架의 Strategy 단계에서 선언 필요"),
    ]
    for i,(icon,title,body) in enumerate(warns):
        add_rect(s, 0.4+i*6.47, 5.1, 6.13, 1.6, fill_color=C_NAVY)
        add_rect(s, 0.4+i*6.47, 5.1, 6.13, 0.05, fill_color=C_GOLD)
        add_textbox(s, f"{icon}  {title}", 0.6+i*6.47, 5.18, 5.7, 0.42,
                    font_size=13, bold=True, color=C_GOLD)
        add_textbox(s, body, 0.6+i*6.47, 5.62, 5.7, 0.95,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)

def slide_16_tools_compare(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "Genspark vs NotebookLM — 언제 무엇을 쓰나?",
                    "두 도구는 경쟁이 아닌 상호 보완 — Promptology 기법으로 극대화")
    draw_part_badge(s, "PART 2 | 연구질문·선행연구")
    draw_bottom_bar(s)
    draw_slide_number(s, 16)

    rows = [
        ("연구 주제 탐색",     "Genspark ✅","—",
         "웹 전체 탐색 · Zero-shot Deep Research"),
        ("연구 질문 초안",     "Genspark ✅","—",
         "Role+CoT 프롬프트로 Gap 분석 극대화"),
        ("확보 문헌 심층 분석","—","NotebookLM ✅",
         "Structured Output으로 표 자동 생성"),
        ("선행연구 매트릭스",  "—","NotebookLM ✅",
         "Few-shot + 인용 추적 기능 활용"),
        ("이론·가설 정교화",   "—","—",
         "Claude + Iterative Refinement"),
        ("통계 코드 생성",     "—","—",
         "Cursor + 바이브 코딩 (Chain-of-Thought)"),
    ]
    headers = ["연구 단계","Genspark","NotebookLM","Promptology 기법"]
    col_ws  = [3.2, 2.3, 2.8, 4.68]
    col_ls  = [0.4, 3.65, 6.0, 8.85]

    for j,(hdr,cw,cl) in enumerate(zip(headers,col_ws,col_ls)):
        bg = C_NAVY if j<3 else C_PURPLE
        add_rect(s, cl, 1.65, cw, 0.52, fill_color=bg)
        add_textbox(s, hdr, cl+0.05, 1.68, cw-0.1, 0.45,
                    font_size=13, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    for i,row in enumerate(rows):
        bg = C_LIGHT_GRAY if i%2==0 else RGBColor(0xE0,0xE8,0xF5)
        for j,(cell,cw,cl) in enumerate(zip(row,col_ws,col_ls)):
            add_rect(s, cl, 2.22+i*0.75, cw, 0.68, fill_color=bg)
            fc = C_BLUE if "✅" in cell else (C_PURPLE if j==3 else C_DARK_GRAY)
            if j==0: fc=C_NAVY
            add_textbox(s, cell, cl+0.05, 2.26+i*0.75, cw-0.1, 0.58,
                        font_size=11, color=fc,
                        align=PP_ALIGN.CENTER if j<3 else PP_ALIGN.LEFT)

def slide_17_model_design_role(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "연구모형 설계 — LLM은 어떤 역할을 할 수 있는가?",
                    "이론적 연결고리를 찾고 변수를 구조화하는 강력한 사고 파트너")
    draw_part_badge(s, "PART 3 | 연구모형 설계")
    draw_bottom_bar(s)
    draw_slide_number(s, 17)

    add_rect(s, 0.4, 1.65, 5.8, 4.5, fill_color=C_TEAL)
    add_rect(s, 0.4, 1.65, 5.8, 0.52, fill_color=RGBColor(0x00,0x4D,0x38))
    add_textbox(s, "✅  LLM이 도울 수 있는 것", 0.4, 1.68, 5.8, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(["이론 탐색 및 비교 분석",
                              "변수 간 관계 정리",
                              "가설 문장 초안 작성",
                              "대안 연구모형 제안",
                              "설문 척도(측정 문항) 초안",
                              "선행연구 기반 모형 정당화"]):
        add_textbox(s, f"  ✦  {item}", 0.55, 2.28+i*0.62, 5.4, 0.52,
                    font_size=12, color=C_WHITE)

    add_rect(s, 7.1, 1.65, 5.8, 4.5, fill_color=C_RED)
    add_rect(s, 7.1, 1.65, 5.8, 0.52, fill_color=RGBColor(0xC0,0x39,0x2B))
    add_textbox(s, "❌  LLM이 대체할 수 없는 것", 7.1, 1.68, 5.8, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(["연구자의 이론적 기여",
                              "맥락적·현장적 판단",
                              "학문적 책임과 윤리",
                              "왜 이 이론인가의 정당화",
                              "가설 방향성 최종 결정",
                              "연구의 독창성과 창의성"]):
        add_textbox(s, f"  ✦  {item}", 7.25, 2.28+i*0.62, 5.4, 0.52,
                    font_size=12, color=C_WHITE)

    add_textbox(s, "⚖️", 6.3, 3.55, 0.7, 0.6,
                font_size=26, align=PP_ALIGN.CENTER)

def slide_18_model_demo(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "🔧 [실습 데모] Claude와 함께 연구모형 설계하기",
                    "SPARRO의 Prompt Design 단계 — CRAFT 모델 적용 실습")
    draw_part_badge(s, "PART 3 | 연구모형 설계")
    draw_bottom_bar(s)
    draw_slide_number(s, 18)

    add_rect(s, 0.4, 1.65, 12.53, 2.1, fill_color=RGBColor(0x1A,0x1A,0x2E))
    add_rect(s, 0.4, 1.65, 0.08, 2.1, fill_color=C_PURPLE)
    add_textbox(s, "💬  CRAFT 모델 적용 프롬프트 (Promptology — Role + CoT + Structured Output)",
                0.6, 1.7, 12.0, 0.38,
                font_size=11, bold=True, color=C_PURPLE)
    add_textbox(s,
                "[C-명확성] 행정학 연구방법론 전문가 역할  "
                "[R-근거] TAM·공공가치론·전자정부 수용모형 비교\n"
                "[A-대상] 행정학 석사 논문 수준  "
                "[F-형식] 변수별 표 + 가설 H1~H5 형식  "
                "[T-과제] 독립·매개·조절·종속변수 포함 모형 + 적합 통계방법 제안",
                0.6, 2.08, 12.1, 1.5,
                font_size=12, color=C_WHITE, word_wrap=True)

    add_rect(s, 0.4, 3.95, 7.2, 2.85,
             fill_color=RGBColor(0x0D,0x0D,0x1A))
    add_textbox(s, "🐍  AI 출력 — 연구모형 및 가설 초안",
                0.6, 4.0, 7.0, 0.38,
                font_size=11, bold=True, color=C_CYAN)
    add_textbox(s,
                "추천 이론: TAM 확장 모형 (행정 맥락 적합)\n\n"
                "독립변수: 서비스 품질(5개 차원)\n"
                "매개변수: 서비스 이용 의도\n"
                "조절변수: 디지털 리터러시\n"
                "종속변수: 주민 만족도\n\n"
                "H1: 서비스 품질 → 만족도 (+)\n"
                "H2: 리터러시 조절효과 (+)\n"
                "추천: SEM + 위계적 회귀",
                0.6, 4.42, 6.8, 2.65,
                font_name="Courier New", font_size=11,
                color=RGBColor(0x00,0xFF,0x7F), word_wrap=True)

    add_rect(s, 7.8, 3.95, 5.1, 2.85, fill_color=C_NAVY)
    add_rect(s, 7.8, 3.95, 5.1, 0.05, fill_color=C_PURPLE)
    add_textbox(s, "📐  SPARRO 검토 포인트",
                8.0, 4.0, 4.8, 0.38,
                font_size=11, bold=True, color=C_PURPLE)
    for i,chk in enumerate(["S: AI 활용 범위 선언 완료?",
                             "P: CRAFT 기법 적용?",
                             "A: 내 논문 맥락과 일치?",
                             "R: 이론 출처 검증 완료?",
                             "R: 가설 방향 타당한가?"]):
        add_textbox(s, f"  {'✅' if i<3 else '🔍'}  {chk}",
                    7.95, 4.45+i*0.47, 4.7, 0.42,
                    font_size=11, color=C_WHITE)

def slide_19_model_vis(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "연구모형 시각화 & 측정 도구 설계",
                    "LLM 출력 → SPARRO Adopting·Refining 단계 → 연구자의 모형 완성")
    draw_part_badge(s, "PART 3 | 연구모형 설계")
    draw_bottom_bar(s)
    draw_slide_number(s, 19)

    boxes = [
        (0.5,  3.1, 2.4, 1.0, "독립변수\n디지털 서비스 품질\n(5개 차원)",     C_BLUE),
        (5.1,  1.8, 2.8, 1.0, "매개변수\n서비스 이용 의도",    C_PURPLE),
        (5.1,  4.3, 2.8, 1.0, "조절변수\n디지털 리터러시",     C_TEAL),
        (10.0, 3.1, 2.6, 1.0, "종속변수\n주민 만족도",         C_RED),
    ]
    for l,t,w,h,label,col in boxes:
        add_rect(s, l, t, w, h, fill_color=col)
        add_textbox(s, label, l+0.05, t+0.05, w-0.1, h-0.1,
                    font_size=12, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    add_textbox(s, "→  H1(+)", 2.95, 3.3, 2.1, 0.5,
                font_size=14, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "→  H2(+)", 7.95, 3.3, 2.0, 0.5,
                font_size=14, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "조절  H3", 6.2, 3.15, 1.5, 0.45,
                font_size=11, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 5.55, 12.53, 0.72, fill_color=C_NAVY)
    add_textbox(s,
                "모형 수식:  Y(만족도) = β₀ + β₁X₁(품질) + β₂X₂(접근성) + β₃M(리터러시) + β₄(X₁×M) + ε",
                0.6, 5.62, 12.1, 0.52,
                font_size=13, bold=True, color=C_CYAN,
                align=PP_ALIGN.CENTER)

    add_textbox(s, "📝  AI 생성 설문 문항 예시 (CRAFT — Structured Output 적용)",
                0.4, 6.4, 7.0, 0.4,
                font_size=12, bold=True, color=C_NAVY)
    add_textbox(s,
                '"이 지방정부 온라인 서비스는 내가 필요한 정보를 정확하게 제공한다 (1=전혀 아님 ~ 5=매우 그렇다)"',
                0.4, 6.78, 12.53, 0.38,
                font_size=11, color=C_DARK_GRAY,
                align=PP_ALIGN.CENTER)

def slide_20_model_limits(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "연구모형 설계 — AI 보조의 한계와 연구자의 역할",
                    "Promptology의 핵심: 프롬프트를 잘 짜는 것 = 연구 설계 능력")
    draw_part_badge(s, "PART 3 | 연구모형 설계")
    draw_bottom_bar(s)
    draw_slide_number(s, 20)

    points = [
        ("🎯","이론 선택의 책임",
         "AI가 여러 이론을 나열해도, '왜 이 이론인가'의 학문적 정당성은 연구자가 설명해야 한다. "
         "SPARRO의 Reviewing 단계에서 반드시 원문 대조 필요."),
        ("↔️","가설 방향성 결정",
         "AI는 양방향 관계를 모두 제시한다. Role 프롬프팅으로 전문가 역할을 부여해도, "
         "방향 결정의 최종 책임은 선행 이론과 연구자 판단에 있다."),
        ("✍️","프롬프트 = 연구 설계 능력",
         "Promptology 관점: CRAFT 모델을 적용한 프롬프트를 작성하려면 "
         "연구 문제를 명확히 알아야 한다. 프롬프트 품질이 곧 연구 설계 품질이다."),
        ("🔍","할루시네이션 검증 필수",
         "AI가 제안한 이론·논문·통계는 반드시 원출처와 대조하라. "
         "SPARRO Reviewing 단계: 존재하지 않는 논문 인용은 빈번히 발생한다."),
    ]
    for i,(icon,title,body) in enumerate(points):
        l = 0.4+(i%2)*6.47
        t = 1.65+(i//2)*2.45
        add_rect(s, l, t, 6.13, 2.2,
                 fill_color=C_NAVY if i%2==0 else C_BLUE)
        add_rect(s, l, t, 6.13, 0.05,
                 fill_color=C_PURPLE if i<2 else C_CYAN)
        add_textbox(s, icon, l+0.1, t+0.1, 0.65, 0.55, font_size=22)
        add_textbox(s, title, l+0.7, t+0.12, 5.2, 0.45,
                    font_size=14, bold=True, color=C_WHITE)
        add_textbox(s, body, l+0.15, t+0.65, 5.8, 1.4,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)

def slide_21_vibe_intro(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "바이브 코딩(Vibe Coding) — 자연어로 코딩하는 시대",
                    "Andrej Karpathy (前 Tesla AI, OpenAI) | 2025년 2월 제시 개념")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 21)

    add_rect(s, 0.4, 1.65, 12.53, 1.1, fill_color=C_NAVY)
    add_rect(s, 0.4, 1.65, 0.1, 1.1, fill_color=C_CYAN)
    add_textbox(s,
                '"코드를 직접 작성하지 않고, 자연어로 AI에게 원하는 것을 말하면 AI가 코드를 생성·실행·디버깅한다"',
                0.65, 1.78, 12.0, 0.82,
                font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 2.95, 5.9, 1.9,
             fill_color=RGBColor(0x6B,0x6B,0x7B))
    add_textbox(s, "❌  기존 패러다임", 0.4, 2.98, 5.9, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "Python / R / SPSS 문법 학습\n→ 코드 작성 → 오류 → 검색 → 수정\n→ 수 시간 ~ 수 일 소요",
                0.55, 3.48, 5.6, 1.2,
                font_size=12, color=C_LIGHT_GRAY, word_wrap=True)

    add_textbox(s, "→", 6.5, 3.65, 0.7, 0.6,
                font_size=28, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    add_rect(s, 7.1, 2.95, 5.8, 1.9, fill_color=C_BLUE)
    add_textbox(s, "✅  바이브 코딩 패러다임", 7.1, 2.98, 5.8, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "자연어로 원하는 것 설명\n→ AI가 코드 즉시 생성·실행\n→ 수 분 내 결과 획득",
                7.25, 3.48, 5.5, 1.2,
                font_size=12, color=C_WHITE, word_wrap=True)

    add_rect(s, 0.4, 5.1, 12.53, 1.6, fill_color=C_TEAL)
    add_textbox(s, "🎓  행정학 연구자에게 의미하는 것",
                0.6, 5.18, 12.0, 0.42,
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(s,
                "Python·R을 몰라도 회귀분석·구조방정식·텍스트 분석 가능  |  "
                "Promptology 핵심: '무엇을 분석할지 아는 것' = 연구 설계 능력이 더 중요\n"
                "🛠️  주요 도구: Cursor · Claude Code · GitHub Copilot · Google AI Studio",
                0.6, 5.65, 12.0, 0.92,
                font_size=12, color=C_WHITE, word_wrap=True)

def slide_22_vibe_compare(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "바이브 코딩 vs 전통적 통계분석 — 무엇이 달라지는가?",
                    "진입 장벽 해소 · 분석 속도 혁신 · 한계는 여전히 존재")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 22)

    headers = ["비교 항목","전통 방식","바이브 코딩"]
    col_ws  = [3.5, 4.2, 4.78]
    col_ls  = [0.4, 3.95, 8.2]
    rows = [
        ("진입 장벽","Python/R/SPSS 학습 필요","자연어 명령만으로 실행 가능"),
        ("오류 대응","스택오버플로우 수동 검색","AI가 즉시 디버깅·수정"),
        ("분석 시간","코드 작성에 수 시간","수 분 내 초안 생성"),
        ("결과 해석","연구자 독립 해석","AI 초안 + 연구자 SPARRO 검토"),
        ("코드 검증","연구자가 직접 확인","⚠️ 반드시 수동 검증 필요"),
        ("통계 적합성","연구자 판단","⚠️ AI는 요청한 분석 무조건 실행"),
    ]

    for j,(hdr,cw,cl) in enumerate(zip(headers,col_ws,col_ls)):
        bg = C_NAVY if j==0 else (C_BLUE if j==1 else C_TEAL)
        add_rect(s, cl, 1.65, cw, 0.52, fill_color=bg)
        add_textbox(s, hdr, cl+0.05, 1.68, cw-0.1, 0.45,
                    font_size=13, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    for i,row in enumerate(rows):
        bg = C_LIGHT_GRAY if i%2==0 else RGBColor(0xE0,0xE8,0xF5)
        warn = "⚠️" in row[1] or "⚠️" in row[2]
        if warn: bg = RGBColor(0xFF,0xF3,0xE0)
        for j,(cell,cw,cl) in enumerate(zip(row,col_ws,col_ls)):
            add_rect(s, cl, 2.22+i*0.75, cw, 0.68, fill_color=bg)
            fc = C_RED if "⚠️" in cell else C_DARK_GRAY
            if j==0: fc=C_NAVY
            add_textbox(s, cell, cl+0.08, 2.26+i*0.75, cw-0.16, 0.58,
                        font_size=12, color=fc)

def slide_23_vibe_demo1(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "🔧 [실습 데모 1] 자연어로 회귀분석 실행하기",
                    "Promptology — Structured Output + Chain-of-Thought 적용")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 23)

    add_rect(s, 0.4, 1.65, 12.53, 2.1, fill_color=RGBColor(0x1A,0x1A,0x2E))
    add_rect(s, 0.4, 1.65, 0.09, 2.1, fill_color=C_PURPLE)
    add_textbox(s, "💬  바이브 코딩 명령 (Promptology: Role + Structured Output + CoT)",
                0.6, 1.7, 12.0, 0.38,
                font_size=11, bold=True, color=C_PURPLE)
    add_textbox(s,
                "[역할] 너는 통계분석 전문가야.  [데이터] survey_data.csv\n"
                "[단계별 실행] ① 기술통계 + 상관분석  "
                "② 다중회귀(종속: satisfaction)  "
                "③ 리터러시 조절효과\n"
                "[형식] APA 양식 표 + Python 코드 + 행정학 맥락 해석문",
                0.6, 2.08, 12.1, 1.5,
                font_size=12, color=C_WHITE, word_wrap=True)

    add_rect(s, 0.4, 3.95, 7.2, 2.85,
             fill_color=RGBColor(0x0D,0x0D,0x1A))
    add_textbox(s, "🐍  AI 생성 Python 코드 (일부)",
                0.6, 4.0, 7.0, 0.38,
                font_size=11, bold=True, color=C_CYAN)
    add_textbox(s,
                "import pandas as pd\n"
                "import statsmodels.api as sm\n\n"
                "df = pd.read_csv('survey_data.csv')\n"
                "print(df.describe())\n\n"
                "X = sm.add_constant(df[['quality','access','literacy']])\n"
                "y = df['satisfaction']\n"
                "model = sm.OLS(y, X).fit()\n"
                "print(model.summary())",
                0.6, 4.42, 6.8, 3.1,
                font_name="Courier New", font_size=11,
                color=RGBColor(0x00,0xFF,0x7F), word_wrap=False)

    add_rect(s, 7.8, 3.95, 5.1, 2.85, fill_color=C_NAVY)
    add_textbox(s, "📊  분석 결과 (요약)",
                8.0, 4.0, 4.8, 0.38,
                font_size=11, bold=True, color=C_GOLD)
    add_textbox(s,
                "R² = 0.542  (p < .001)\n\n"
                "quality:   β=0.42  p<.001  ***\n"
                "access:    β=0.21  p<.01   **\n"
                "literacy:  β=0.18  p<.05   *\n\n"
                "조절효과 (quality×literacy):\n"
                "β=0.15  p<.05  *",
                7.95, 4.42, 4.7, 2.65,
                font_name="Courier New", font_size=11,
                color=C_WHITE, word_wrap=False)

def slide_24_vibe_demo2(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "🔧 [실습 데모 2] 통계 결과를 AI로 해석하기",
                    "SPARRO Reviewing·Refining 단계 — AI 초안 → 연구자 검토·완성")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 24)

    add_rect(s, 0.4, 1.65, 12.53, 1.3, fill_color=RGBColor(0x1A,0x1A,0x2E))
    add_rect(s, 0.4, 1.65, 0.09, 1.3, fill_color=C_PURPLE)
    add_textbox(s, "💬  Iterative Refinement 프롬프트 (SPARRO Reviewing 적용)",
                0.6, 1.7, 12.0, 0.38,
                font_size=11, bold=True, color=C_PURPLE)
    add_textbox(s,
                '"아래 회귀분석 결과를 행정학 석사 논문 결과 섹션에 맞게 해석해줘. '
                'APA 형식 + 통계적 유의성 + H1~H3 연결 포함. [결과 붙여넣기]"',
                0.6, 2.05, 12.1, 0.75,
                font_size=12, color=C_WHITE, word_wrap=True)

    add_rect(s, 0.4, 3.15, 12.53, 3.1, fill_color=C_NAVY)
    add_rect(s, 0.4, 3.15, 12.53, 0.05, fill_color=C_GOLD)
    add_textbox(s, "📝  AI 생성 해석 초안 (SPARRO Adopting → Reviewing → Refining 필요)",
                0.6, 3.22, 12.0, 0.42,
                font_size=12, bold=True, color=C_GOLD)
    add_textbox(s,
                "디지털 서비스 품질(β = .42, p < .001)은 주민 만족도에 유의한 정(+)의 영향을 미쳤으며, H1을 지지한다.\n\n"
                "디지털 리터러시의 조절효과(β = .15, p < .05)가 유의하여 H3을 지지하였다. "
                "이는 리터러시가 높은 주민일수록 서비스 품질의 긍정적 영향이 더 크게 나타남을 의미하며, "
                "TAM 이론의 맥락과 일치한다.\n\n"
                "전체 모형 설명력은 R² = .542로 종속변수 분산의 54.2%를 설명하였다.",
                0.6, 3.7, 12.1, 2.35,
                font_size=12, color=C_WHITE, word_wrap=True)

    add_textbox(s,
                "⚠️  SPARRO Reviewing 필수: AI 해석은 초안 — 맥락 오류·과대해석 가능 | "
                "Refining 단계에서 연구자가 이론적 맥락 보완",
                0.4, 6.38, 12.53, 0.42,
                font_size=11, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

def slide_25_vibe_limits(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "바이브 코딩의 가능성과 한계 — 연구자가 반드시 알아야 할 것",
                    "도구의 힘을 최대화하되 맹목적 의존은 경계한다")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 25)

    poss = ["코딩 장벽 제거 → 더 많은 연구자가 정량 연구 접근 가능",
            "분석 속도 향상 → 더 다양한 모형 시도 가능",
            "결과 해석 초안 자동화 → 반복 서술 작업 절약",
            "오류 디버깅 AI 자동화 → 방법론 실험 용이"]
    limits = ["생성 코드 반드시 수동 검증 (오류도 자신있게 출력)",
              "통계 방법 적합성은 연구자 판단 (AI는 무조건 실행)",
              "결과 해석 과도한 의존 금지 — 통계 의미 이해 필수",
              "민감 데이터 보안 — 개인정보 외부 서버 업로드 주의"]

    add_rect(s, 0.4, 1.65, 5.9, 4.85, fill_color=RGBColor(0x00,0x5C,0x3A))
    add_rect(s, 0.4, 1.65, 5.9, 0.52, fill_color=RGBColor(0x00,0x4D,0x30))
    add_textbox(s, "✅  가능성", 0.4, 1.68, 5.9, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(poss):
        add_textbox(s, f"  ✦  {item}", 0.55, 2.28+i*0.88, 5.6, 0.75,
                    font_size=12, color=C_WHITE, word_wrap=True)

    add_rect(s, 7.0, 1.65, 5.9, 4.85, fill_color=RGBColor(0x7B,0x1F,0x1F))
    add_rect(s, 7.0, 1.65, 5.9, 0.52, fill_color=RGBColor(0x5C,0x15,0x15))
    add_textbox(s, "⚠️  한계", 7.0, 1.68, 5.9, 0.45,
                font_size=14, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    for i,item in enumerate(limits):
        add_textbox(s, f"  ⚠  {item}", 7.15, 2.28+i*0.88, 5.6, 0.75,
                    font_size=12, color=C_WHITE, word_wrap=True)

    add_textbox(s, "⚖️", 6.3, 3.8, 0.7, 0.65,
                font_size=26, align=PP_ALIGN.CENTER)

def slide_26_vibe_analytics(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "Vibe Analytics — 바이브 코딩의 다음 단계",
                    "MIT Sloan Management Review, 2025.7 | 데이터 분석의 민주화")
    draw_part_badge(s, "PART 4 | 통계분석")
    draw_bottom_bar(s)
    draw_slide_number(s, 26)

    add_rect(s, 0.4, 1.65, 12.53, 0.88, fill_color=C_NAVY)
    add_textbox(s,
                '"바이브 코딩이 프로그래밍에 하는 것을, 바이브 애널리틱스는 데이터 분석에 한다"',
                0.6, 1.75, 12.1, 0.65,
                font_size=16, bold=True, color=C_CYAN,
                align=PP_ALIGN.CENTER)

    flow = [
        ("📂","데이터 업로드","공공 데이터\ne-나라지표\n통계청 KOSIS"),
        ("🔍","자동 탐색 분석","AI가 스스로\n패턴 발견\n이상치 탐지"),
        ("📈","시각화 자동","그래프·차트\n자동 생성\n인터랙티브"),
        ("💡","인사이트 도출","정책 트렌드\n자동 분석\n해석 초안"),
    ]
    for i,(icon,title,body) in enumerate(flow):
        l = 0.4+i*3.15
        add_rect(s, l, 2.75, 2.8, 3.1,
                 fill_color=C_BLUE if i%2==0 else C_NAVY)
        add_rect(s, l, 2.75, 2.8, 0.06,
                 fill_color=C_CYAN if i%2==0 else C_GOLD)
        add_textbox(s, icon, l, 2.85, 2.8, 0.65,
                    font_size=26, align=PP_ALIGN.CENTER)
        add_textbox(s, title, l, 3.48, 2.8, 0.55,
                    font_size=14, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, body, l+0.1, 4.1, 2.6, 1.5,
                    font_size=12,
                    color=C_CYAN if i%2==0 else C_LIGHT_GRAY,
                    align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(s, "▶", l+2.8+0.13, 4.1, 0.3, 0.6,
                        font_size=16, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 6.05, 12.53, 0.72, fill_color=C_TEAL)
    add_textbox(s,
                "🎯  연구자의 새 역할: 분석을 '하는' 사람 → 분석을 '설계하고 판단하는' 사람 "
                "| Promptology = 이 역할을 위한 핵심 역량",
                0.6, 6.15, 12.1, 0.52,
                font_size=13, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

def slide_27_workflow(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "AI 보조 연구 전 과정 — 통합 워크플로우",
                    "Promptology가 모든 단계를 연결하는 핵심 역량이다")
    draw_part_badge(s, "PART 5 | 종합")
    draw_bottom_bar(s)
    draw_slide_number(s, 27)

    pipeline = [
        ("①","연구 아이디어","Genspark\nZero-shot",C_BLUE),
        ("②","연구 질문(RQ)","Genspark\nRole+CoT",RGBColor(0x0D,0x47,0xA1)),
        ("③","선행연구 검토","NotebookLM\nStructured Out",C_TEAL),
        ("④","연구모형 설계","Claude\nCRAFT모델",C_PURPLE),
        ("⑤","데이터 수집","설문 플랫폼\n공공 데이터",RGBColor(0x7B,0x3F,0x00)),
        ("⑥","통계분석","바이브 코딩\nCoT 프롬프트",C_RED),
        ("⑦","해석·작성","SPARRO\nIterative",RGBColor(0x1A,0x5C,0x2A)),
    ]
    box_w, box_h, start_l, top = 1.75, 2.85, 0.3, 1.75
    for i,(num,step,tool,col) in enumerate(pipeline):
        l = start_l+i*1.83
        add_rect(s, l, top, box_w, box_h, fill_color=col)
        add_rect(s, l, top, box_w, 0.05, fill_color=C_CYAN)
        add_textbox(s, num, l, top+0.08, box_w, 0.45,
                    font_size=18, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, step, l, top+0.5, box_w, 0.75,
                    font_size=12, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_rect(s, l+0.2, top+1.32, box_w-0.4, 0.04, fill_color=C_CYAN)
        add_textbox(s, tool, l, top+1.44, box_w, 0.92,
                    font_size=10, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        if i < 6:
            add_textbox(s, "▶", l+box_w+0.02, top+1.1, 0.3, 0.45,
                        font_size=12, color=C_GOLD,
                        align=PP_ALIGN.CENTER)

    # Promptology 강조 배너
    add_rect(s, 0.3, 4.82, 12.73, 0.55, fill_color=C_PURPLE)
    add_textbox(s,
                "✍️  Promptology: Zero-shot → Few-shot → Role → CoT → Structured Output → Iterative Refinement",
                0.5, 4.9, 12.4, 0.4,
                font_size=12, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    add_rect(s, 0.3, 5.55, 12.73, 0.72, fill_color=C_NAVY)
    add_textbox(s,
                "⏱  전통적 연구 준비: 3~6개월  →  AI + Promptology 보조: 2~4주 (약 4~8배 단축)",
                0.5, 5.65, 12.4, 0.52,
                font_size=14, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    principles = ["① Human-in-the-Loop","② 출처 검증 필수",
                  "③ 투명한 보고","④ 연구 아젠다 선도"]
    add_rect(s, 0.3, 6.38, 12.73, 0.75, fill_color=RGBColor(0x1A,0x1A,0x2E))
    for i,p in enumerate(principles):
        l = 0.5+i*3.15
        add_textbox(s, p, l, 6.5, 3.0, 0.5,
                    font_size=12, bold=True, color=C_CYAN,
                    align=PP_ALIGN.CENTER)

def slide_28_principles_final(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "행정학 연구자를 위한 AI 활용 4대 원칙",
                    "기술을 넘어 — 학문적 책임과 연구자의 정체성")
    draw_part_badge(s, "PART 5 | 종합")
    draw_bottom_bar(s)
    draw_slide_number(s, 28)

    principles = [
        ("①","Augmentation,\nnot Replacement",
         "AI는 연구자를 대체하는 것이 아니라 능력을 증폭시킨다. "
         "Promptology는 이 증폭의 언어이다.",C_BLUE),
        ("②","비판적 AI 리터러시",
         "AI 출력을 항상 의심하고 검증하라. SPARRO Reviewing 단계를 "
         "절대 건너뛰지 마라.",C_RED),
        ("③","투명한 보고",
         "AI 활용 도구·단계·프롬프트를 논문 방법론에 명시하라. "
         "COREQ+LLM(2025) 기준 준수.",C_TEAL),
        ("④","연구 아젠다 선도",
         "AI 거버넌스·책임성·공정성 연구를 선도하라. "
         "Promptology는 행정학의 새 방법론 언어다.",C_PURPLE),
    ]
    for i,(num,title,body,col) in enumerate(principles):
        l = 0.4+(i%2)*6.47
        t = 1.65+(i//2)*2.5
        add_rect(s, l, t, 6.13, 2.25, fill_color=col)
        add_rect(s, l, t, 6.13, 0.05,
                 fill_color=C_CYAN if i%2==0 else C_GOLD)
        add_rect(s, l, t, 0.62, 2.25, fill_color=RGBColor(0,0,0))
        add_textbox(s, num, l, t, 0.62, 2.25,
                    font_size=26, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title, l+0.72, t+0.12, 5.2, 0.72,
                    font_size=14, bold=True, color=C_WHITE)
        add_textbox(s, body, l+0.72, t+0.85, 5.2, 1.25,
                    font_size=11, color=C_LIGHT_GRAY, word_wrap=True)

def slide_29_closing(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "AI 시대 행정학 연구자의 새로운 역량",
                    "오늘부터 당장 시작할 수 있는 실천 3가지")
    draw_part_badge(s, "PART 5 | 종합")
    draw_bottom_bar(s)
    draw_slide_number(s, 29)

    add_rect(s, 0.4, 1.65, 12.53, 0.85, fill_color=C_NAVY)
    add_textbox(s,
                '"20년 전 행정학자에게 SPSS는 낯선 도구였다. 지금은 기본 소양이다.'
                ' Promptology·LLM·AI Agent도 마찬가지다."',
                0.6, 1.75, 12.1, 0.65,
                font_size=14, bold=True, color=C_CYAN,
                align=PP_ALIGN.CENTER)

    actions = [
        ("🔍","내일 당장","Genspark에서 자신의\n연구 주제로\nDeep Research 실행",C_BLUE),
        ("📚","이번 학기","논문 5편을 NotebookLM에\n업로드 + Structured\nOutput 질문 실습",C_TEAL),
        ("✍️","이번 주","CRAFT 모델로\n연구 관련 프롬프트 3개\n직접 설계해보기",C_PURPLE),
    ]
    for i,(icon,when,action,col) in enumerate(actions):
        l = 0.7+i*4.1
        add_rect(s, l, 2.72, 3.6, 3.6, fill_color=col)
        add_rect(s, l, 2.72, 3.6, 0.06,
                 fill_color=C_CYAN if i<2 else C_GOLD)
        add_textbox(s, icon, l, 2.82, 3.6, 0.72,
                    font_size=30, align=PP_ALIGN.CENTER)
        add_textbox(s, when, l, 3.52, 3.6, 0.45,
                    font_size=14, bold=True,
                    color=C_GOLD if i==0 else C_CYAN,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, action, l+0.1, 4.05, 3.4, 2.0,
                    font_size=13, color=C_WHITE,
                    align=PP_ALIGN.CENTER, word_wrap=True)

    add_rect(s, 0.4, 6.52, 12.53, 0.35, fill_color=C_GOLD)
    add_textbox(s,
                "🏆  Promptology를 아는 연구자가 AI 시대의 행정학을 이끈다",
                0.6, 6.55, 12.1, 0.28,
                font_size=13, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

def slide_30_references(prs):
    s = blank_slide(prs)
    draw_header_bar(s, "참고 자료 & 도구 링크 | Q&A",
                    "오늘 강의에서 소개한 도구·Promptology 문헌·프레임워크")
    draw_bottom_bar(s)
    draw_slide_number(s, 30)

    add_textbox(s, "🛠️  주요 도구", 0.4, 1.65, 6.0, 0.42,
                font_size=14, bold=True, color=C_NAVY)
    tools = [
        ("Genspark",    "genspark.ai",             C_BLUE),
        ("NotebookLM",  "notebooklm.google.com",   C_TEAL),
        ("Claude",      "claude.ai",               C_PURPLE),
        ("Cursor",      "cursor.com",              RGBColor(0x7B,0x3F,0x00)),
        ("Claude Code", "docs.anthropic.com/claude-code", C_RED),
        ("Google Colab","colab.research.google.com",RGBColor(0x1A,0x5C,0x2A)),
    ]
    for i,(name,url,col) in enumerate(tools):
        t = 2.18+i*0.68
        add_rect(s, 0.4, t, 5.9, 0.58, fill_color=col)
        add_textbox(s, name, 0.55, t+0.07, 2.0, 0.42,
                    font_size=12, bold=True, color=C_WHITE)
        add_textbox(s, url,  2.55, t+0.07, 3.5, 0.42,
                    font_size=11, color=C_CYAN)

    add_textbox(s, "📚  주요 참고 문헌", 6.7, 1.65, 6.5, 0.42,
                font_size=14, bold=True, color=C_NAVY)
    refs = [
        "Olla et al., Promptology: Enhancing Human-AI Interaction\n(MDPI Information, 2024)",
        "Scaling Hermeneutics: qualitative coding with LLMs\n(EPJ Data Science, 2025)",
        "COREQ+LLM: Reporting LLM use in qualitative research\n(Research Protocols, 2025)",
        "Vibe Analytics: Vibe Coding's New Cousin\n(MIT Sloan Management Review, 2025)",
        "NIA, AI 활용 정부서비스 사례와 유형 분석 (2025)",
        "공공부문 AI 에이전트 도입 및 활용 현황 (정부혁신, 2025)",
    ]
    for i,ref in enumerate(refs):
        bg = C_LIGHT_GRAY if i%2==0 else RGBColor(0xE0,0xE8,0xF5)
        t = 2.18+i*0.73
        add_rect(s, 6.7, t, 6.28, 0.65, fill_color=bg)
        add_textbox(s, ref, 6.85, t+0.04, 5.98, 0.58,
                    font_size=10, color=C_DARK_GRAY, word_wrap=True)

    add_rect(s, 0.4, 6.62, 12.53, 0.62, fill_color=C_NAVY)
    add_textbox(s, "💬  Q & A  —  질문이 있으신가요?  |  강의자료 공유 예정",
                0.6, 6.7, 12.1, 0.45,
                font_size=15, bold=True, color=C_CYAN,
                align=PP_ALIGN.CENTER)

# ── 4. 전체 실행 ─────────────────────────────────────────────

def build_presentation():
    prs = new_prs()

    print("=" * 55)
    print("  행정학 연구 AI LLM & Promptology 강의 PPT 생성 중")
    print("=" * 55)

    funcs = [
        (slide_01_cover,                  "표지"),
        (slide_02_roadmap,                "강의 로드맵"),
        (slide_03_why,                    "왜 배워야 하는가"),
        (slide_04_llm_vs_agent,           "LLM vs AI Agent"),
        (slide_05_tools,                  "AI 도구 5종"),
        (slide_06_principles,             "AI 활용 원칙 3가지"),
        (slide_07_promptology_def,        "★ Promptology 개념 정의"),
        (slide_08_promptology_techniques, "★ Promptology 핵심 기법 6가지"),
        (slide_09_sparro,                 "★ SPARRO 프레임워크"),
        (slide_10_rq_design,              "연구 질문 설계"),
        (slide_11_genspark_demo,          "Genspark 데모"),
        (slide_12_finer,                  "FINER 기준"),
        (slide_13_notebooklm,             "NotebookLM 소개"),
        (slide_14_notebooklm_demo,        "NotebookLM 데모"),
        (slide_15_methodology,            "방법론적 의미"),
        (slide_16_tools_compare,          "Genspark vs NotebookLM"),
        (slide_17_model_design_role,      "연구모형 설계 역할"),
        (slide_18_model_demo,             "연구모형 설계 데모"),
        (slide_19_model_vis,              "연구모형 시각화"),
        (slide_20_model_limits,           "연구모형 AI 한계"),
        (slide_21_vibe_intro,             "바이브 코딩 개념"),
        (slide_22_vibe_compare,           "바이브 코딩 비교"),
        (slide_23_vibe_demo1,             "바이브 코딩 데모1"),
        (slide_24_vibe_demo2,             "결과 해석 데모"),
        (slide_25_vibe_limits,            "바이브 코딩 한계"),
        (slide_26_vibe_analytics,         "Vibe Analytics"),
        (slide_27_workflow,               "통합 워크플로우"),
        (slide_28_principles_final,       "4대 원칙 최종"),
        (slide_29_closing,                "마무리"),
        (slide_30_references,             "참고자료 & Q&A"),
    ]

    for i, (func, name) in enumerate(funcs, 1):
        func(prs)
        print(f"  [{i:02d}/30] Slide {i:02d} — {name}  ✓")

    fname = "AI_행정학연구_Promptology_강의.pptx"
    prs.save(fname)
    print()
    print("=" * 55)
    print(f"  ✅ 완료!  '{fname}'  생성 완료")
    print("  📥 Google Colab: 왼쪽 파일 탐색기에서 다운로드")
    print("  📁 로컬 환경: 현재 폴더에서 확인")
    print("=" * 55)
    return fname

# ── 실행 ─────────────────────────────────────────────────────
build_presentation()

# Google Colab 자동 다운로드 (선택)
try:
    from google.colab import files
    files.download("AI_행정학연구_Promptology_강의.pptx")
    print("📥 다운로드가 시작되었습니다.")
except ImportError:
    print("📁 로컬 환경: 현재 폴더에서 파일을 확인하세요.")